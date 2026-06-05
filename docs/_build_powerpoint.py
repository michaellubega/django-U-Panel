#!/usr/bin/env python3
"""
Build U_Panel_KIU_Overview.pptx from PROJECT_OVERVIEW.html.

Default (native + screenshots):
  - Headings and paragraphs → real PowerPoint text (Calibri, colors from :root).
  - Each diagram: **Inkscape → EMF** when available; else **Playwright PNG** of the
    rendered `figure.diagram` (browser-accurate); else a short placeholder.
  - Also writes **U_Panel_KIU_Overview_screenshots.pptx** — full-slide PNG deck (unless
    --native-only).

Optional:
  --screenshot   Only the full-slide PNG deck (writes U_Panel_KIU_Overview_screenshots.pptx).
  --native-only  Skip the full-slide screenshot deck.

Inkscape (optional, vector diagrams): https://inkscape.org/

Prerequisites:
  pip install -r docs/requirements-powerpoint.txt
  playwright install chromium
"""

from __future__ import annotations

import argparse
import io
import re
import shutil
import struct
import subprocess
import sys
import tempfile
from pathlib import Path

from bs4 import BeautifulSoup, NavigableString


def _blank_layout(prs):
    for lo in prs.slide_layouts:
        if lo.name and str(lo.name).lower() == "blank":
            return lo
    return prs.slide_layouts[6]


def _parse_css_colors(style_text: str) -> dict[str, tuple[int, int, int]]:
    """Extract #rrggbb from :root { --name: #hex; }"""
    out: dict[str, tuple[int, int, int]] = {}
    for m in re.finditer(r"--([\w-]+)\s*:\s*#([0-9a-fA-F]{6})", style_text):
        name, hx = m.group(1), m.group(2)
        r, g, b = int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16)
        out[name] = (r, g, b)
    return out


def _find_inkscape() -> str | None:
    for name in ("inkscape", "inkscape.com"):
        p = shutil.which(name)
        if p:
            return p
    win = Path(r"C:\Program Files\Inkscape\bin\inkscape.com")
    if win.is_file():
        return str(win)
    win2 = Path(r"C:\Program Files\Inkscape\bin\inkscape.exe")
    if win2.is_file():
        return str(win2)
    return None


def _uniquify_svg_ids(svg_xml: str, prefix: str) -> str:
    soup = BeautifulSoup(svg_xml, "xml")
    root = soup.find("svg")
    if root is None:
        return svg_xml
    id_map: dict[str, str] = {}
    for tag in root.find_all(True):
        if not tag.has_attr("id"):
            continue
        old = tag["id"]
        if old in id_map:
            continue
        new = f"{prefix}{old}"
        id_map[old] = new
        tag["id"] = new
    raw = str(root)
    for old, new in id_map.items():
        raw = raw.replace(f"url(#{old})", f"url(#{new})")
        raw = raw.replace(f'url("#{old}")', f'url("#{new}")')
    return raw


def _svg_viewbox_size(svg_xml: str) -> tuple[float, float]:
    m = re.search(r'viewBox\s*=\s*["\']\s*([\d.\-\s]+)\s*["\']', svg_xml)
    if not m:
        return 16.0, 9.0
    parts = [float(x) for x in m.group(1).replace(",", " ").split() if x.strip()]
    if len(parts) >= 4:
        w, h = parts[2], parts[3]
        if w > 0 and h > 0:
            return w, h
    return 16.0, 9.0


def _svg_to_emf(inkscape: str, svg_content: str, emf_path: Path) -> bool:
    svg_path = emf_path.with_suffix(".svg")
    svg_path.write_text(svg_content, encoding="utf-8")
    try:
        r = subprocess.run(
            [
                inkscape,
                str(svg_path),
                "--export-type=emf",
                f"--export-filename={emf_path}",
                "--export-background=white",
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
        return r.returncode == 0 and emf_path.is_file()
    except (OSError, subprocess.TimeoutExpired):
        return False


def _png_size(data: bytes) -> tuple[int, int]:
    if len(data) < 24 or data[:8] != b"\x89PNG\r\n\x1a\n":
        raise ValueError("not a PNG")
    w = int.from_bytes(data[16:20], "big")
    h = int.from_bytes(data[20:24], "big")
    return w, h


class DiagramScreenshots:
    """Lazy Playwright session to PNG-capture figure.diagram elements in document order."""

    def __init__(self, html_path: Path):
        self.html_path = html_path.resolve()
        self._pw = None
        self._browser = None
        self._page = None

    def _open(self) -> None:
        if self._page is not None:
            return
        from playwright.sync_api import sync_playwright

        self._pw = sync_playwright().start()
        self._browser = self._pw.chromium.launch(headless=True)
        self._page = self._browser.new_page(
            device_scale_factor=2,
            viewport={"width": 1920, "height": 1080},
        )
        self._page.goto(self.html_path.as_uri(), wait_until="networkidle", timeout=120_000)
        try:
            self._page.evaluate("() => document.fonts.ready")
        except Exception:
            pass
        self._page.wait_for_timeout(500)

    def capture_diagram(self, zero_based_index: int) -> bytes | None:
        try:
            self._open()
        except Exception:
            return None
        try:
            loc = self._page.locator("figure.diagram").nth(zero_based_index)
            loc.wait_for(state="visible", timeout=20_000)
            return loc.screenshot(type="png")
        except Exception:
            return None

    def close(self) -> None:
        if self._browser is not None:
            self._browser.close()
            self._browser = None
        if self._pw is not None:
            self._pw.stop()
            self._pw = None
        self._page = None


def _init_overview_ppt_js() -> str:
    return r"""() => {
  if (window.__pptShow) return;
  const root = document.querySelector(".page");
  if (!root) return;
  const hero = root.querySelector(":scope > header.hero");
  const main = root.querySelector(":scope > main");
  const footer = root.querySelector(":scope > footer");
  const mainChildren = main ? Array.prototype.slice.call(main.children) : [];
  const groups = [];
  let cur = [];
  for (let i = 0; i < mainChildren.length; i++) {
    const el = mainChildren[i];
    if (el.tagName === "H2" && cur.length) {
      groups.push(cur);
      cur = [el];
    } else {
      cur.push(el);
    }
  }
  if (cur.length) groups.push(cur);
  window.__pptSlideCount = 1 + groups.length;
  window.__pptShow = function (n) {
    const G = groups;
    const last = G.length;
    if (hero) hero.style.display = n === 0 ? "" : "none";
    if (footer) footer.style.display = n === last ? "" : "none";
    mainChildren.forEach(function (el) { el.style.display = "none"; });
    if (n > 0 && n <= last) {
      const gi = n - 1;
      G[gi].forEach(function (el) { el.style.display = ""; });
    }
  };
}"""


def _viewport_height(page) -> int:
    h = page.evaluate(
        """() => {
      const el = document.documentElement;
      const b = document.body;
      return Math.ceil(Math.max(el.scrollHeight, b ? b.scrollHeight : 0, el.clientHeight));
    }"""
    )
    if not isinstance(h, int) or h < 1:
        return 1080
    return min(12000, max(1080, h + 48))


def build_screenshot_pptx(html: Path, out: Path) -> int:
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Emu, Inches

    vw = 1920
    uri = html.as_uri()
    shots: list[bytes] = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(
            device_scale_factor=2,
            viewport={"width": vw, "height": 1080},
        )
        page.goto(uri, wait_until="networkidle", timeout=60_000)
        try:
            page.evaluate("() => document.fonts.ready")
        except Exception:
            pass
        page.wait_for_timeout(400)
        page.evaluate(_init_overview_ppt_js())
        n = page.evaluate("() => window.__pptSlideCount")
        if not isinstance(n, int) or n < 1:
            print("Could not parse overview into slides.", file=sys.stderr)
            browser.close()
            return 1
        for i in range(n):
            page.evaluate("(idx) => window.__pptShow(idx)", i)
            page.wait_for_timeout(250)
            vh = _viewport_height(page)
            page.set_viewport_size({"width": vw, "height": vh})
            page.wait_for_timeout(200)
            page.evaluate("() => window.scrollTo(0, 0)")
            page.wait_for_timeout(100)
            shots.append(page.screenshot(type="png", full_page=False))
        browser.close()

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = _blank_layout(prs)
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)

    for png in shots:
        iw, ih = _png_size(png)
        slide = prs.slides.add_slide(blank)
        scale = min(sw / iw, sh / ih)
        w = int(scale * iw)
        h = int(scale * ih)
        left = int((sw - w) / 2)
        top = int((sh - h) / 2)
        slide.shapes.add_picture(
            io.BytesIO(png),
            Emu(left),
            Emu(top),
            width=Emu(w),
            height=Emu(h),
        )

    try:
        prs.save(out)
    except PermissionError:
        print("Could not write (close PowerPoint if open):", out, file=sys.stderr)
        return 1
    print("Wrote", out, "slides", len(shots), "(screenshot mode)")
    return 0


def _rgb(cols: dict[str, tuple[int, int, int]], key: str, default: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    t = cols.get(key, default)
    return RGBColor(*t)


def _add_text_block(slide, left, top, width, height, text, *, pt, bold=False, rgb=None):
    from pptx.util import Emu, Pt

    if not (text or "").strip():
        return top
    le = int(left)
    te = int(top)
    we = int(width)
    he = int(height)
    box = slide.shapes.add_textbox(le, te, we, he)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text.strip()
    p.font.size = Pt(pt)
    p.font.bold = bold
    if rgb is not None:
        p.font.color.rgb = rgb
    p.font.name = "Calibri"
    return Emu(te + he)


def _render_block(
    slide,
    el,
    top,
    left,
    content_w,
    inkscape: str | None,
    svg_idx: list[int],
    cols: dict[str, tuple[int, int, int]],
    diagram_shots: DiagramScreenshots | None,
):
    from pptx.util import Emu, Inches

    text_rgb = _rgb(cols, "text", (10, 41, 21))
    brand_dark = _rgb(cols, "brand-dark", (15, 77, 46))
    muted = _rgb(cols, "muted", (100, 116, 139))

    name = getattr(el, "name", None)
    if name is None:
        return top

    if name == "h2":
        h = Inches(0.38)
        return _add_text_block(
            slide,
            left,
            top,
            content_w,
            h,
            el.get_text(" ", strip=True),
            pt=20,
            bold=True,
            rgb=brand_dark,
        )

    if name == "p":
        cls = " ".join(el.get("class", []) or [])
        pt = 13 if "tight" in cls or "lede" in cls else 14
        h = Inches(0.85)
        return _add_text_block(
            slide, left, top, content_w, h, el.get_text(" ", strip=True), pt=pt, rgb=text_rgb
        )

    if name == "figure":
        cap_el = el.find("figcaption")
        svg_el = el.find("svg")
        if svg_el is None:
            return top
        svg_idx[0] += 1
        prefix = f"d{svg_idx[0]}_"
        raw = _uniquify_svg_ids(str(svg_el), prefix)
        vw, vh = _svg_viewbox_size(raw)
        aspect = vw / max(vh, 1e-6)
        cw = int(content_w)
        top_emu = int(top)
        slide_h = int(Inches(7.5))
        max_h = slide_h - top_emu - int(Inches(0.35))
        pic_w = cw
        pic_h = int(pic_w * vh / vw)
        min_h = int(Inches(0.6))
        if pic_h > max_h and max_h > min_h:
            pic_h = max_h
            pic_w = int(pic_h * vw / vh)
        emf_ok = False
        with tempfile.TemporaryDirectory() as td:
            tdir = Path(td)
            emf = tdir / f"fig{svg_idx[0]}.emf"
            if inkscape and _svg_to_emf(inkscape, raw, emf):
                slide.shapes.add_picture(
                    str(emf), int(left), top_emu, width=pic_w, height=pic_h
                )
                emf_ok = True
        if emf_ok:
            top = Emu(top_emu + pic_h + int(Inches(0.12)))
        else:
            png_ok = False
            if diagram_shots is not None:
                png = diagram_shots.capture_diagram(svg_idx[0] - 1)
                if png:
                    try:
                        slide.shapes.add_picture(
                            io.BytesIO(png),
                            int(left),
                            top_emu,
                            width=pic_w,
                            height=pic_h,
                        )
                        png_ok = True
                        top = Emu(top_emu + pic_h + int(Inches(0.12)))
                    except Exception:
                        png_ok = False
            if not png_ok:
                note = (
                    "[Diagram omitted - install Inkscape (EMF) or Playwright+Chromium "
                    "(PNG snapshots). See docs/requirements-powerpoint.txt]"
                )
                top = _add_text_block(
                    slide,
                    left,
                    top,
                    content_w,
                    Inches(0.55),
                    note,
                    pt=11,
                    rgb=muted,
                )
        if cap_el:
            cap = cap_el.get_text(" ", strip=True)
            if cap:
                top = _add_text_block(
                    slide,
                    left,
                    top,
                    content_w,
                    Inches(0.35),
                    cap,
                    pt=10,
                    rgb=muted,
                )
        return Emu(int(top) + int(Inches(0.15)))

    if name == "ul":
        lines = []
        for li in el.find_all("li", recursive=False):
            t = li.get_text(" ", strip=True)
            if t:
                lines.append("• " + t)
        if lines:
            h = Inches(min(2.8, 0.32 * len(lines) + 0.2))
            return _add_text_block(
                slide, left, top, content_w, h, "\n".join(lines), pt=13, rgb=text_rgb
            )
        return top

    return top


def _hero_slide(prs, hero, cols):
    from pptx.util import Emu, Inches

    blank = _blank_layout(prs)
    slide = prs.slides.add_slide(blank)
    left = Emu(int(Inches(0.45)))
    top = Emu(int(Inches(0.4)))
    content_w = Emu(int(Inches(12.4)))

    for ch in hero.children:
        if isinstance(ch, NavigableString) or not getattr(ch, "name", None):
            continue
        if ch.name == "h1":
            title = ch.get_text(" ", strip=True)
            top = _add_text_block(
                slide,
                left,
                top,
                content_w,
                Inches(0.75),
                title,
                pt=36,
                bold=True,
                rgb=_rgb(cols, "brand-dark", (15, 77, 46)),
            )
        elif ch.name == "p":
            cls = " ".join(ch.get("class", []) or [])
            pt = 12 if "meta" in cls else 15
            muted = "meta" in cls
            top = _add_text_block(
                slide,
                left,
                top,
                content_w,
                Inches(0.55),
                ch.get_text(" ", strip=True),
                pt=pt,
                rgb=_rgb(cols, "muted", (100, 116, 139)) if muted else _rgb(cols, "text", (10, 41, 21)),
            )


def _main_groups(main) -> list[list]:
    groups: list[list] = []
    cur: list = []
    for child in main.children:
        if isinstance(child, NavigableString):
            continue
        if not getattr(child, "name", None):
            continue
        if child.name == "h2" and cur:
            groups.append(cur)
            cur = [child]
        else:
            cur.append(child)
    if cur:
        groups.append(cur)
    return groups


def build_native_pptx(
    html: Path,
    out: Path,
    diagram_shots: DiagramScreenshots | None,
) -> tuple[int, str]:
    from pptx import Presentation
    from pptx.util import Emu, Inches

    raw_html = html.read_text(encoding="utf-8")
    cols = _parse_css_colors(raw_html)
    soup = BeautifulSoup(raw_html, "lxml")
    page = soup.select_one(".page")
    if page is None:
        print("Missing .page in HTML.", file=sys.stderr)
        return 1, "error"
    hero = page.select_one("header.hero")
    main = page.select_one("main")
    footer = page.select_one("footer")
    if hero is None or main is None:
        print("Missing header.hero or main.", file=sys.stderr)
        return 1, "error"

    inkscape = _find_inkscape()
    if inkscape is None and diagram_shots is None:
        print(
            "Note: No Inkscape and no Playwright - diagram placeholders only.",
            file=sys.stderr,
        )
    elif inkscape is None:
        print("Note: Inkscape not found - diagrams use Playwright PNG snapshots.", file=sys.stderr)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    svg_counter = [0]

    _hero_slide(prs, hero, cols)

    groups = _main_groups(main)
    last_gi = len(groups) - 1

    for gi, group in enumerate(groups):
        slide = prs.slides.add_slide(_blank_layout(prs))
        left = Emu(int(Inches(0.45)))
        top = Emu(int(Inches(0.35)))
        content_w = Emu(int(Inches(12.4)))
        for node in group:
            top = _render_block(
                slide,
                node,
                top,
                left,
                content_w,
                inkscape,
                svg_counter,
                cols,
                diagram_shots,
            )
        if gi == last_gi and footer:
            ft = footer.get_text(" ", strip=True)
            if ft:
                top = Emu(int(top) + int(Inches(0.35)))
                _add_text_block(
                    slide,
                    left,
                    top,
                    content_w,
                    Inches(0.45),
                    ft,
                    pt=12,
                    rgb=_rgb(cols, "muted", (100, 116, 139)),
                )

    if inkscape:
        mode = "emf"
    elif diagram_shots is not None:
        mode = "png"
    else:
        mode = "none"
    try:
        prs.save(out)
    except PermissionError:
        print("Could not write (close PowerPoint if open):", out, file=sys.stderr)
        return 1, "error"
    print(
        "Wrote",
        out,
        "- native text +",
        "vector EMF (Inkscape)" if mode == "emf" else "PNG diagrams (Playwright)" if mode == "png" else "diagram placeholders",
    )
    return 0, mode


def main() -> int:
    ap = argparse.ArgumentParser(description="Build PowerPoint from PROJECT_OVERVIEW.html")
    ap.add_argument(
        "--screenshot",
        action="store_true",
        help="Only the full-slide PNG deck (U_Panel_KIU_Overview_screenshots.pptx).",
    )
    ap.add_argument(
        "--native-only",
        action="store_true",
        help="Skip writing U_Panel_KIU_Overview_screenshots.pptx after the native deck.",
    )
    args = ap.parse_args()

    docs = Path(__file__).resolve().parent
    html = docs / "PROJECT_OVERVIEW.html"
    out_native = docs / "U_Panel_KIU_Overview.pptx"
    out_shots = docs / "U_Panel_KIU_Overview_screenshots.pptx"
    if not html.is_file():
        print("Missing:", html, file=sys.stderr)
        return 1

    if args.screenshot:
        try:
            import playwright.sync_api  # noqa: F401
        except ImportError:
            print("Install: pip install playwright && playwright install chromium", file=sys.stderr)
            return 1
        return build_screenshot_pptx(html, out_shots)

    try:
        import bs4  # noqa: F401
    except ImportError:
        print("Install: pip install beautifulsoup4 lxml", file=sys.stderr)
        return 1
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("Install: pip install python-pptx", file=sys.stderr)
        return 1

    diagram_shots: DiagramScreenshots | None = None
    try:
        import playwright.sync_api  # noqa: F401

        diagram_shots = DiagramScreenshots(html)
    except ImportError:
        pass

    try:
        rc, _mode = build_native_pptx(html, out_native, diagram_shots)
    finally:
        if diagram_shots is not None:
            diagram_shots.close()

    if rc != 0:
        return rc

    if args.native_only:
        return 0

    try:
        import playwright.sync_api  # noqa: F401
    except ImportError:
        print("Skipping slide screenshot deck (install playwright for U_Panel_KIU_Overview_screenshots.pptx).")
        return 0

    r2 = build_screenshot_pptx(html, out_shots)
    if r2 != 0:
        print("Native deck OK; slide screenshot deck failed (see errors above).", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
